"""Extracts plain text from an uploaded document's raw bytes.

Format-dependent, but everything downstream (chunking, embedding) is
format-agnostic - this is the one step whose whole job is normalizing messy
real-world files down to one common shape.

Supported: plain text, PDF (text-layer via pypdf, per-page - any page with
no text layer, or whose text layer turns out to be unusable (see
_has_pua_glyphs), falls back to OCR just for that page, see
_extract_pdf_text), Word (.docx), PowerPoint (.pptx), Excel (.xlsx) - modern
Office Open XML formats only, not the legacy binary .doc/.ppt/.xls formats
(those need much heavier tooling, e.g. a LibreOffice conversion step -
deliberately out of scope for now).
"""

import asyncio
import io
import logging

import openpyxl
import pytesseract
from docx import Document as DocxDocument
from pdf2image import convert_from_bytes
from pptx import Presentation
from pypdf import PdfReader

logger = logging.getLogger(__name__)

# DPI for OCR rasterization - high enough for decent recognition accuracy on
# typical scans without ballooning render/OCR time on large documents.
_OCR_DPI = 300

# BMP Private Use Area (U+E000-U+F8FF). Some PDFs - math/logic worksheets
# out of an old Word Equation Editor or MathType are the recurring case -
# embed their symbol font's glyphs (minus, times, tau, ...) mapped into
# this range instead of the real Unicode math characters. pypdf extracts
# exactly what the font's CMap says, so the "text" it recovers for those
# glyphs is technically valid Unicode but meaningless everywhere else -
# nothing downstream (this app, a browser, a terminal) has that font's
# private mapping table, so those characters render as nothing and a
# formula like "p ∧ q → r" comes back as "p q r" with the operators just
# gone. extract_text() doesn't return empty for a page like this - it
# returns real prose *plus* invisible junk - so it isn't caught by the
# "no text layer" check below and needs its own detection.
_PRIVATE_USE_AREA = range(0xE000, 0xF900)


def _has_pua_glyphs(text: str) -> bool:
    return any(ord(ch) in _PRIVATE_USE_AREA for ch in text)


class UnsupportedDocumentError(Exception):
    pass


_DOCX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_XLSX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


async def extract_text(content_type: str, data: bytes) -> str:
    if content_type == "text/plain":
        return _require_nonempty(data.decode("utf-8", errors="replace"))

    if content_type == "application/pdf":
        return await _extract_pdf_text(data)

    if content_type == _DOCX_CONTENT_TYPE:
        return _extract_docx_text(data)

    if content_type == _PPTX_CONTENT_TYPE:
        return _extract_pptx_text(data)

    if content_type == _XLSX_CONTENT_TYPE:
        return _extract_xlsx_text(data)

    raise UnsupportedDocumentError(f"unsupported content type: {content_type}")


def _require_nonempty(text: str) -> str:
    if not text.strip():
        raise UnsupportedDocumentError("no extractable text found")
    return text


async def _extract_pdf_text(data: bytes) -> str:
    """Per-page, not per-document: a page either has a native text layer, or
    it doesn't, and that's decided independently for each page - a mixed
    document (some real text pages, some scanned/image pages) is common
    (e.g. a native PDF with one scanned cover or diagram page), and deciding
    OCR at the whole-document level would either skip OCR entirely (if *any*
    page has text, the old check on the joined text of every page passed)
    or waste it on every page (if none does). Neither matches what's
    actually needed here.

    No attempt to pre-classify *why* a page has no text (image vs. blank vs.
    broken) before deciding whether to OCR it - some PDFs draw "text" as
    vector line art rather than real font-embedded text objects, which
    would look like neither to a naive images-present check but would still
    render to legible pixels OCR can read. Simpler and more robust to just
    try OCR on every page with no native text and let the outcome answer
    the question: recovers something -> it was an image page; still empty
    after OCR -> genuinely blank or broken, logged and skipped, not fatal
    on its own.
    """
    reader = PdfReader(io.BytesIO(data))
    page_texts: list[str] = []
    needs_ocr: list[int] = []  # 0-indexed positions into page_texts

    for i, page in enumerate(reader.pages):
        try:
            extracted = page.extract_text() or ""
        except Exception:
            # A corrupt content stream on one page shouldn't take down
            # extraction for every other (possibly perfectly fine) page in
            # the same document.
            extracted = ""
        page_texts.append(extracted)
        if not extracted.strip() or _has_pua_glyphs(extracted):
            needs_ocr.append(i)

    tesseract_missing = False
    if needs_ocr:
        tesseract_available = True
        for i in needs_ocr:
            if not tesseract_available:
                break
            try:
                ocr_text = await asyncio.to_thread(_ocr_page_sync, data, i + 1)  # pdf2image is 1-indexed
            except pytesseract.TesseractNotFoundError:
                # A deploy/config problem (Tesseract isn't installed on this
                # host), not a per-page document-quality one - stop
                # attempting further pages (every one would fail the same
                # way) rather than retrying the same broken setup N times.
                tesseract_available = False
                tesseract_missing = True
                logger.warning(
                    "extractor.ocr_unavailable",
                    extra={"event": "extractor.ocr_unavailable", "reason": "tesseract binary not installed"},
                )
                continue

            if ocr_text.strip():
                # Also the right outcome for a PUA-glyph page (not just a
                # genuinely empty one): OCR reads the rendered pixels, not
                # the font's internal codes, so it recovers the real symbol
                # a PUA-mapped font hid from native extraction.
                page_texts[i] = ocr_text
                logger.info(
                    "extractor.page_recovered_via_ocr",
                    extra={"event": "extractor.page_recovered_via_ocr", "page": i + 1},
                )
            elif page_texts[i].strip():
                # Had native text (this was a PUA-glyph page, not an empty
                # one - see _has_pua_glyphs), but OCR couldn't recover
                # anything better. Keep the native text rather than
                # discarding it - prose with some invisible junk mixed in
                # still beats nothing - but it's still degraded, worth a log.
                logger.warning(
                    "extractor.page_pua_glyphs_ocr_failed",
                    extra={"event": "extractor.page_pua_glyphs_ocr_failed", "page": i + 1},
                )
            else:
                # No native text and OCR found nothing either - genuinely
                # blank or broken. Not fatal on its own; only an error if
                # it turns out every page in the document is like this.
                logger.warning(
                    "extractor.page_has_no_content",
                    extra={"event": "extractor.page_has_no_content", "page": i + 1},
                )

    text = "\n\n".join(p for p in page_texts if p.strip())
    if not text.strip():
        if tesseract_missing:
            raise UnsupportedDocumentError("OCR unavailable - the Tesseract binary isn't installed on this host")
        raise UnsupportedDocumentError(
            "no extractable text found on any page, even after OCR - the file may be blank, "
            "corrupt, or too low-quality to read"
        )
    return text


def _ocr_page_sync(data: bytes, page_number: int) -> str:
    """Renders and OCRs a single page (1-indexed, pdf2image's own
    convention). One page at a time, not the whole document up front -
    a full-page render at _OCR_DPI is several MB uncompressed; rendering
    every page before OCR touches any of them is what made large scanned
    documents run out of memory in the same process serving the rest of the
    API. grayscale halves the per-page memory footprint vs RGB with no OCR
    accuracy cost - Tesseract works off grayscale internally regardless of
    the input's color mode.
    """
    images = convert_from_bytes(data, dpi=_OCR_DPI, first_page=page_number, last_page=page_number, grayscale=True)
    if not images:
        return ""
    return pytesseract.image_to_string(images[0])


def _extract_docx_text(data: bytes) -> str:
    doc = DocxDocument(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                parts.append(row_text)
    return _require_nonempty("\n\n".join(parts))


def _extract_pptx_text(data: bytes) -> str:
    presentation = Presentation(io.BytesIO(data))
    slides_text = []
    for i, slide in enumerate(presentation.slides, start=1):
        shape_texts = [
            shape.text_frame.text for shape in slide.shapes if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if shape_texts:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(shape_texts))
    return _require_nonempty("\n\n".join(slides_text))


def _extract_xlsx_text(data: bytes) -> str:
    workbook = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    sheets_text = []
    for sheet in workbook.worksheets:
        rows_text = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows_text.append(" | ".join(cells))
        if rows_text:
            sheets_text.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows_text))
    return _require_nonempty("\n\n".join(sheets_text))
