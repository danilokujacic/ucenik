from unittest.mock import AsyncMock, patch

from tests.conftest import auth_headers, fake_completion, login
from ucenik.enum.user_role import UserRole


async def _create_subject(client, headers, name="Biology") -> str:
    response = await client.post("/subjects", json={"name": name}, headers=headers)
    assert response.status_code == 201, response.text
    return response.json()["id"]


async def test_download_returns_the_original_bytes(client, make_user):
    await make_user("teacher@x.com", UserRole.TEACHER)
    student = await make_user("student@x.com", UserRole.STUDENT)
    t_tokens = await login(client, "teacher@x.com", "password123")
    s_tokens = await login(client, "student@x.com", "password123")
    subject_id = await _create_subject(client, auth_headers(t_tokens))

    enroll = await client.post(
        f"/subjects/{subject_id}/enrollments", json={"student_id": str(student.id)}, headers=auth_headers(t_tokens)
    )
    assert enroll.status_code == 201

    content = b"unique content for the download test"
    with patch("ucenik.services.documents.ingest_document", AsyncMock()):
        upload = await client.post(
            f"/subjects/{subject_id}/documents",
            files={"file": ("notes.txt", content, "text/plain")},
            headers=auth_headers(t_tokens),
        )
    document_id = upload.json()["id"]

    # an enrolled student (not just the owning teacher) can download too -
    # same access level as reading document metadata (require_subject_access)
    response = await client.get(
        f"/subjects/{subject_id}/documents/{document_id}/download", headers=auth_headers(s_tokens)
    )

    assert response.status_code == 200
    assert response.content == content
    assert response.headers["content-type"].startswith("text/plain")
    assert 'filename="notes.txt"' in response.headers["content-disposition"]


async def test_download_available_even_before_ingest_finishes(client, make_user):
    """The raw file lands in storage synchronously on upload, before ingest
    even runs - download shouldn't be gated on ingest status.
    """
    await make_user("teacher@x.com", UserRole.TEACHER)
    tokens = await login(client, "teacher@x.com", "password123")
    subject_id = await _create_subject(client, auth_headers(tokens))

    with patch("ucenik.services.documents.ingest_document", AsyncMock()):
        upload = await client.post(
            f"/subjects/{subject_id}/documents",
            files={"file": ("notes.txt", b"still pending", "text/plain")},
            headers=auth_headers(tokens),
        )
    assert upload.json()["status"] == "pending"
    document_id = upload.json()["id"]

    response = await client.get(
        f"/subjects/{subject_id}/documents/{document_id}/download", headers=auth_headers(tokens)
    )

    assert response.status_code == 200
    assert response.content == b"still pending"


async def test_outsider_cannot_download(client, make_user):
    await make_user("teacher@x.com", UserRole.TEACHER)
    await make_user("outsider@x.com", UserRole.STUDENT)
    t_tokens = await login(client, "teacher@x.com", "password123")
    outsider_tokens = await login(client, "outsider@x.com", "password123")
    subject_id = await _create_subject(client, auth_headers(t_tokens))

    with patch("ucenik.services.documents.ingest_document", AsyncMock()):
        upload = await client.post(
            f"/subjects/{subject_id}/documents",
            files={"file": ("notes.txt", b"secret notes", "text/plain")},
            headers=auth_headers(t_tokens),
        )
    document_id = upload.json()["id"]

    response = await client.get(
        f"/subjects/{subject_id}/documents/{document_id}/download", headers=auth_headers(outsider_tokens)
    )

    assert response.status_code == 403


async def test_student_cannot_upload_document(client, make_user):
    await make_user("teacher@x.com", UserRole.TEACHER)
    await make_user("student@x.com", UserRole.STUDENT)
    t_tokens = await login(client, "teacher@x.com", "password123")
    s_tokens = await login(client, "student@x.com", "password123")
    subject_id = await _create_subject(client, auth_headers(t_tokens))

    response = await client.post(
        f"/subjects/{subject_id}/documents",
        files={"file": ("notes.txt", b"hello world", "text/plain")},
        headers=auth_headers(s_tokens),
    )

    assert response.status_code == 403


async def test_unsupported_content_type_rejected(client, make_user):
    await make_user("teacher@x.com", UserRole.TEACHER)
    tokens = await login(client, "teacher@x.com", "password123")
    subject_id = await _create_subject(client, auth_headers(tokens))

    with patch("ucenik.services.documents.ingest_document", AsyncMock()):
        response = await client.post(
            f"/subjects/{subject_id}/documents",
            files={"file": ("image.png", b"\x89PNG\r\n", "image/png")},
            headers=auth_headers(tokens),
        )

    assert response.status_code == 415


async def test_upload_creates_document_and_dedups_on_reupload(client, make_user):
    await make_user("teacher@x.com", UserRole.TEACHER)
    tokens = await login(client, "teacher@x.com", "password123")
    subject_id = await _create_subject(client, auth_headers(tokens))

    with patch("ucenik.services.documents.ingest_document", AsyncMock()) as mock_ingest:
        first = await client.post(
            f"/subjects/{subject_id}/documents",
            files={"file": ("notes.txt", b"unique content for dedup test", "text/plain")},
            headers=auth_headers(tokens),
        )
        assert first.status_code == 201
        assert first.json()["status"] == "pending"
        assert mock_ingest.call_count == 1

        # re-uploading identical content should be a no-op, not a new document
        second = await client.post(
            f"/subjects/{subject_id}/documents",
            files={"file": ("notes-renamed.txt", b"unique content for dedup test", "text/plain")},
            headers=auth_headers(tokens),
        )
        assert second.status_code == 201
        assert second.json()["id"] == first.json()["id"]
        assert mock_ingest.call_count == 1  # not called again


async def test_list_and_get_document_requires_subject_access(client, make_user):
    await make_user("teacher@x.com", UserRole.TEACHER)
    await make_user("outsider@x.com", UserRole.STUDENT)
    t_tokens = await login(client, "teacher@x.com", "password123")
    outsider_tokens = await login(client, "outsider@x.com", "password123")
    subject_id = await _create_subject(client, auth_headers(t_tokens))

    with patch("ucenik.services.documents.ingest_document", AsyncMock()):
        upload = await client.post(
            f"/subjects/{subject_id}/documents",
            files={"file": ("notes.txt", b"list and get test content", "text/plain")},
            headers=auth_headers(t_tokens),
        )
    document_id = upload.json()["id"]

    listing = await client.get(f"/subjects/{subject_id}/documents", headers=auth_headers(t_tokens))
    assert listing.status_code == 200
    assert len(listing.json()) == 1

    detail = await client.get(f"/subjects/{subject_id}/documents/{document_id}", headers=auth_headers(t_tokens))
    assert detail.status_code == 200

    # not enrolled, not the owner -> 403
    forbidden = await client.get(f"/subjects/{subject_id}/documents", headers=auth_headers(outsider_tokens))
    assert forbidden.status_code == 403


async def test_delete_document_removes_it(client, make_user):
    await make_user("teacher@x.com", UserRole.TEACHER)
    tokens = await login(client, "teacher@x.com", "password123")
    subject_id = await _create_subject(client, auth_headers(tokens))

    with patch("ucenik.services.documents.ingest_document", AsyncMock()):
        upload = await client.post(
            f"/subjects/{subject_id}/documents",
            files={"file": ("notes.txt", b"delete me please unique content", "text/plain")},
            headers=auth_headers(tokens),
        )
    document_id = upload.json()["id"]

    delete = await client.delete(f"/subjects/{subject_id}/documents/{document_id}", headers=auth_headers(tokens))
    assert delete.status_code == 204

    get_after = await client.get(f"/subjects/{subject_id}/documents/{document_id}", headers=auth_headers(tokens))
    assert get_after.status_code == 404


async def test_full_ingest_pipeline_end_to_end(client, make_user):
    """Real extraction, chunking, embedding, caching, and Chroma storage -
    mocks only the LLM proxy call (real proxy exists now, src/ucenik/llm_proxy/,
    but hitting it here would mean a real Groq API call - see
    tests/test_llm_proxy.py for coverage of the proxy service itself).
    """
    await make_user("teacher@x.com", UserRole.TEACHER)
    tokens = await login(client, "teacher@x.com", "password123")
    subject_id = await _create_subject(client, auth_headers(tokens))

    text = (
        b"Mitosis is the process by which a single cell divides into two identical "
        b"daughter cells. It occurs in four main stages: prophase, metaphase, "
        b"anaphase, and telophase."
    )

    with patch(
        "ucenik.rag.contextualizer.complete",
        AsyncMock(return_value=fake_completion("This chunk is from a biology chapter on mitosis.")),
    ):
        upload = await client.post(
            f"/subjects/{subject_id}/documents",
            files={"file": ("mitosis.txt", text, "text/plain")},
            headers=auth_headers(tokens),
        )
        assert upload.status_code == 201
        document_id = upload.json()["id"]

        detail = await client.get(f"/subjects/{subject_id}/documents/{document_id}", headers=auth_headers(tokens))

    assert detail.status_code == 200
    body = detail.json()
    assert body["status"] == "ready", body
    assert body["chunk_count"] >= 1

    from ucenik.rag.embedder import embed_query
    from ucenik.rag.vector_store import query_similar

    query_vector = await embed_query("what are the stages of mitosis?")
    results = await query_similar(subject_id, query_vector, n_results=1)
    assert "mitosis" in results["documents"][0][0].lower()


async def test_docx_upload_ingests_successfully(client, make_user):
    import io

    from docx import Document as DocxDocument

    await make_user("teacher@x.com", UserRole.TEACHER)
    tokens = await login(client, "teacher@x.com", "password123")
    subject_id = await _create_subject(client, auth_headers(tokens))

    doc = DocxDocument()
    doc.add_paragraph("Mitosis is the process of cell division into two daughter cells.")
    buf = io.BytesIO()
    doc.save(buf)

    with patch(
        "ucenik.rag.contextualizer.complete", AsyncMock(return_value=fake_completion("A biology note on mitosis."))
    ):
        upload = await client.post(
            f"/subjects/{subject_id}/documents",
            files={
                "file": (
                    "notes.docx",
                    buf.getvalue(),
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            },
            headers=auth_headers(tokens),
        )
        assert upload.status_code == 201
        detail = await client.get(
            f"/subjects/{subject_id}/documents/{upload.json()['id']}", headers=auth_headers(tokens)
        )

    assert detail.json()["status"] == "ready", detail.json()


async def test_pptx_upload_ingests_successfully(client, make_user):
    import io

    from pptx import Presentation

    await make_user("teacher@x.com", UserRole.TEACHER)
    tokens = await login(client, "teacher@x.com", "password123")
    subject_id = await _create_subject(client, auth_headers(tokens))

    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[1])
    slide.shapes.title.text = "Mitosis: an overview of cell division"
    buf = io.BytesIO()
    pres.save(buf)

    with patch("ucenik.rag.contextualizer.complete", AsyncMock(return_value=fake_completion("A slide on mitosis."))):
        upload = await client.post(
            f"/subjects/{subject_id}/documents",
            files={
                "file": (
                    "slides.pptx",
                    buf.getvalue(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
            headers=auth_headers(tokens),
        )
        assert upload.status_code == 201
        detail = await client.get(
            f"/subjects/{subject_id}/documents/{upload.json()['id']}", headers=auth_headers(tokens)
        )

    assert detail.json()["status"] == "ready", detail.json()


async def test_xlsx_upload_ingests_successfully(client, make_user):
    import io

    import openpyxl

    await make_user("teacher@x.com", UserRole.TEACHER)
    tokens = await login(client, "teacher@x.com", "password123")
    subject_id = await _create_subject(client, auth_headers(tokens))

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Phase", "Description"])
    ws.append(["Prophase", "Chromatin condenses into chromosomes during mitosis"])
    buf = io.BytesIO()
    wb.save(buf)

    with patch(
        "ucenik.rag.contextualizer.complete",
        AsyncMock(return_value=fake_completion("A spreadsheet on mitosis phases.")),
    ):
        upload = await client.post(
            f"/subjects/{subject_id}/documents",
            files={
                "file": (
                    "phases.xlsx",
                    buf.getvalue(),
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
            },
            headers=auth_headers(tokens),
        )
        assert upload.status_code == 201
        detail = await client.get(
            f"/subjects/{subject_id}/documents/{upload.json()['id']}", headers=auth_headers(tokens)
        )

    assert detail.json()["status"] == "ready", detail.json()


async def test_filename_is_sanitized(client, make_user):
    await make_user("teacher@x.com", UserRole.TEACHER)
    tokens = await login(client, "teacher@x.com", "password123")
    subject_id = await _create_subject(client, auth_headers(tokens))

    with patch("ucenik.services.documents.ingest_document", AsyncMock()):
        upload = await client.post(
            f"/subjects/{subject_id}/documents",
            files={"file": ("<script>alert(1)</script>.txt", b"harmless unique content here", "text/plain")},
            headers=auth_headers(tokens),
        )

    assert upload.status_code == 201
    filename = upload.json()["filename"]
    assert "<" not in filename
    assert ">" not in filename
    assert "/" not in filename


async def test_oversized_upload_rejected_via_content_length_before_reading(client, make_user):
    """The Content-Length pre-check (api/documents.py) should reject before
    ever calling file.read() - patching MAX_FILE_SIZE down to something tiny
    makes this fast to exercise without actually uploading tens of MB.
    """
    await make_user("teacher@x.com", UserRole.TEACHER)
    tokens = await login(client, "teacher@x.com", "password123")
    subject_id = await _create_subject(client, auth_headers(tokens))

    with (
        patch("ucenik.api.documents.MAX_FILE_SIZE", 10),
        patch("ucenik.services.documents.ingest_document", AsyncMock()),
    ):
        response = await client.post(
            f"/subjects/{subject_id}/documents",
            files={"file": ("big.txt", b"this is well over ten bytes", "text/plain")},
            headers=auth_headers(tokens),
        )

    assert response.status_code == 413
